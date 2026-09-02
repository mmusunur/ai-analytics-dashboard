"""
Generate architecture & workflow presentation for AI Analytics Dashboard.
Output: AI_Analytics_Dashboard_Presentation.pptx (project root, alongside tasks.md)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
OUTPUT = ROOT / "AI_Analytics_Dashboard_Presentation.pptx"

# Theme — dark analytics dashboard palette
BG_DARK = RGBColor(15, 23, 42)       # slate-900
BG_CARD = RGBColor(30, 27, 75)       # indigo-950
ACCENT = RGBColor(124, 58, 237)      # violet-600
ACCENT_LIGHT = RGBColor(167, 139, 250)
TEXT_PRIMARY = RGBColor(226, 232, 240)
TEXT_MUTED = RGBColor(148, 163, 184)
SUCCESS = RGBColor(52, 211, 153)
WARNING = RGBColor(251, 191, 36)
WHITE = RGBColor(255, 255, 255)


def set_slide_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.33), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_footer(slide, text="AgenticOps AI · AI Analytics Dashboard"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER


def add_title_block(slide, title: str, subtitle: str = ""):
    add_accent_bar(slide)
    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(1.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT_LIGHT
        p2.space_before = Pt(8)


def add_bullets(slide, items: list[tuple[str, str]], left=0.7, top=1.5, width=11.8, font_size=15):
    """items: list of (bold_prefix, body) or plain strings."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            prefix, body = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run_b = p.add_run()
            run_b.text = prefix
            run_b.font.bold = True
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = ACCENT_LIGHT
            run_t = p.add_run()
            run_t.text = body
            run_t.font.size = Pt(font_size)
            run_t.font.color.rgb = TEXT_PRIMARY
        else:
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = TEXT_PRIMARY
            p.space_after = Pt(10)
        if i > 0:
            tf.paragraphs[-1].space_before = Pt(6)


def add_card(slide, left, top, width, height, title, lines: list[str], accent=ACCENT):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_CARD
    rect.line.color.rgb = accent
    rect.line.width = Pt(1.5)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = accent
    for line in lines:
        pl = tf.add_paragraph()
        pl.text = line
        pl.font.size = Pt(11)
        pl.font.color.rgb = TEXT_PRIMARY
        pl.space_before = Pt(4)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ─────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_slide_background(s1, BG_DARK)
    tb = s1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "AI Analytics Dashboard"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Autonomous Agent Platform — Architecture & End-to-End Workflow"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_LIGHT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "React + FastAPI + MCP Agents + Plane PM"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)
    add_footer(s1)

    # ── Slide 2: End-to-End Workflow ───────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_slide_background(s2, BG_DARK)
    add_title_block(s2, "End-to-End Autonomous Workflow", "From Plane task to Done — zero manual steps after To Do")
    steps = [
        ("1. Pickup  ", "Sprint Watcher polls Plane → moves task to In Progress"),
        ("2. Build  ", "Builder Agent implements code (NLP intent → React/FastAPI changes)"),
        ("3. Test  ", "Tester Agent runs unit + browser + sprint acceptance cases"),
        ("4. Close  ", "Plane Agent marks task Completed when quality gate passes"),
        ("5. Git Push  ", "Git Agent commits meaningful file changes only"),
        ("6. Done  ", "UI pipeline tracker + Task Queue show ✓ Done"),
    ]
    add_bullets(s2, steps, top=1.45, font_size=14)
    # Flow arrow hint
    flow = s2.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(12), Inches(0.4))
    fp = flow.text_frame.paragraphs[0]
    fp.text = "Human creates task in Plane (To Do once)  →  Agent handles everything else automatically"
    fp.font.size = Pt(12)
    fp.font.color.rgb = SUCCESS
    fp.alignment = PP_ALIGN.CENTER
    add_footer(s2)

    # ── Slide 3: Pipeline Gates & Smart Performance ────────────────────────
    s3 = prs.slides.add_slide(blank)
    set_slide_background(s3, BG_DARK)
    add_title_block(s3, "Quality Gates & Smart Test Modes", "Task 37 — no step skipping; fast verify-close")
    add_card(
        s3, Inches(0.5), Inches(1.45), Inches(5.9), Inches(2.4),
        "Mandatory Gates",
        [
            "Build must pass before Test runs",
            "Test failure → return to Build (not Close)",
            "Checkmarks only when step actually passed",
            "Auto-retry — no manual To Do moves",
        ],
        ACCENT,
    )
    add_card(
        s3, Inches(6.7), Inches(1.45), Inches(5.9), Inches(2.4),
        "Smart Performance",
        [
            "Full tests after code change (~5–15 min)",
            "Fast verify-close: smoke + sprint cases (~1 min)",
            "Active poll: 15s · Idle poll: 30s",
            "Test heartbeat visible in Sprint Board UI",
        ],
        SUCCESS,
    )
    add_card(
        s3, Inches(0.5), Inches(4.1), Inches(12.1), Inches(2.3),
        "Live Telemetry (memory/agent_state.json)",
        [
            "pipeline.phase · progress_pct · test_subphase · completed_steps",
            "task_queue: pending / active / completed / failed",
            "agent_working flag — UI pauses polling only during Build (file edits)",
        ],
        WARNING,
    )
    add_footer(s3)

    # ── Slide 4: FastAPI (dedicated) ───────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    set_slide_background(s5, BG_DARK)
    add_title_block(s5, "FastAPI Backend", "Python 3.10 · backend/main.py · OpenAPI at /docs")
    add_card(
        s5, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.0),
        "Core Routers",
        [
            "/api/data — CSV/Excel upload, DB queries",
            "/api/analytics — AI Copilot (no-date NL search)",
            "/api/charts — bar, scatter, KPI aggregates",
            "/api/warehouse/statistics — filtered table + pagination",
            "/api/sprints — Plane tasks, agent status, watcher control",
            "/api/mcp — MCP server registry & health",
            "/api/agents/status — pipeline + task queue telemetry",
            "/api/health — uptime probe for tests & monitor",
        ],
        RGBColor(59, 130, 246),
    )
    add_card(
        s5, Inches(6.6), Inches(1.4), Inches(5.9), Inches(5.0),
        "Design Highlights",
        [
            "CORS enabled for Vite frontend",
            "Dual search: Global Header (date+DB+whse) vs Copilot (no date)",
            "Parameter propagation: oerdte, target_db, oewhse → all widgets",
            "Sprint API caches Plane responses (rate-limit safe)",
            "Warehouse service: live DB + local seed fallback",
            "Serves Agent Monitor UI with live JSON state",
        ],
        ACCENT,
    )
    add_footer(s5)

    # ── Slide 6: MCP (dedicated) ───────────────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    set_slide_background(s6, BG_DARK)
    add_title_block(s6, "Model Context Protocol (MCP)", "Standardized tool bridge — agents/mcp_client.py + mcp_config.json")
    mcp_servers = [
        ("Plane MCP", "list_tasks · update_task_status · add_comment · list_projects", "Sprint Watcher, Plane Agent"),
        ("GitHub MCP", "git_commit · git_push · get_meaningful_changed_files", "Git Agent after task close"),
        ("Memory MCP", "load_state · set_pipeline_status · get_task_queue · log_task_result", "All agents + UI"),
        ("Browser MCP", "run_unit_tests · run_browser_tests · run_sprint_task_tests", "Tester Agent quality gate"),
    ]
    x_positions = [0.5, 6.7]
    for idx, (name, tools, used) in enumerate(mcp_servers):
        col = idx % 2
        row = idx // 2
        add_card(
            s6,
            Inches(x_positions[col]),
            Inches(1.45 + row * 2.55),
            Inches(5.9),
            Inches(2.35),
            name,
            [f"Tools: {tools}", f"Used by: {used}"],
            SUCCESS if idx % 2 else ACCENT_LIGHT,
        )
    note = s6.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(12), Inches(0.45))
    np = note.text_frame.paragraphs[0]
    np.text = "MCP decouples agents from integrations — swap Plane/GitHub/test backends without changing UI or watcher logic"
    np.font.size = Pt(11)
    np.font.color.rgb = TEXT_MUTED
    np.alignment = PP_ALIGN.CENTER
    add_footer(s6)

    # ── Slide 7: Data Flow ─────────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    set_slide_background(s7, BG_DARK)
    add_title_block(s7, "Data Flow — Dashboard to Database", "Single-warehouse filtering · dual search modes")
    add_bullets(
        s7,
        [
            "Global Header Submit → oerdte + target_db + oewhse → /api/charts/* + /api/warehouse/statistics",
            "AI Copilot Ask → POST /api/analytics/ai-copilot with oerdte='' (all dates, NL intent)",
            "KPI cards, bar chart, scatter plot, warehouse table refresh from same filter context",
            "When warehouse selected: charts show ONLY that facility (Task 27 rule)",
            "Frontend polls /api/agents/status for pipeline progress (10s interval, pauses on Build only)",
        ],
        top=1.5,
        font_size=14,
    )
    add_footer(s7)

    # ── Slide 7: Summary ───────────────────────────────────────────────────
    s8 = prs.slides.add_slide(blank)
    set_slide_background(s8, BG_DARK)
    tb = s8.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    bullets = [
        "Autonomous 6-step pipeline: Plane → Build → Test → Close → Git → Done",
        "FastAPI powers analytics, warehouse data, and agent telemetry APIs",
        "MCP connects agents to Plane, Git, memory state, and Playwright tests",
        "Smart fast/full test modes keep quality high while closing tasks quickly",
    ]
    for i, b in enumerate(bullets):
        pl = tf.add_paragraph()
        pl.text = f"✓  {b}"
        pl.font.size = Pt(15)
        pl.font.color.rgb = TEXT_PRIMARY if i else TEXT_PRIMARY
        pl.space_before = Pt(14)
        pl.alignment = PP_ALIGN.LEFT
    add_footer(s8, "Questions? · See tasks.md & /docs for full specification")

    prs.save(str(OUTPUT))
    return prs


if __name__ == "__main__":
    build_presentation()
    print(f"Created: {OUTPUT}")
