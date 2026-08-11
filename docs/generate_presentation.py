"""
Generate AgenticOps AI overview presentation (PowerPoint).
Run: python docs/generate_presentation.py
Output: docs/AgenticOps_AI_Overview.pptx

Content synced from docs/doc_content.py — run docs/sync_all_documentation.py to regenerate all.
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as exc:
    raise SystemExit("Install python-pptx: pip install python-pptx") from exc

from doc_content import PPTX_SLIDES
from doc_outputs import PPTX_PATH, remove_legacy_pptx_copies, save_in_place

OUT = PPTX_PATH


def _build_presentation(dest: Path):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for title, body in PPTX_SLIDES:
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True

        bx = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(9), Inches(5))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(18)

    prs.save(str(dest))


def main():
    remove_legacy_pptx_copies()
    save_in_place(OUT, _build_presentation)
    print(f"[OK] Presentation replaced in place: {OUT}")


if __name__ == "__main__":
    main()
